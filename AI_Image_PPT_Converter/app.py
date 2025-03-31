import os
import base64
import io
import tempfile
import re
from flask import Flask, request, jsonify, send_file, send_from_directory
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE, MSO_CONNECTOR
from pptx.enum.dml import MSO_LINE, MSO_FILL
from PIL import Image
import openai

app = Flask(__name__, static_folder='static')

# Configuration
openai.api_key = os.getenv('OPENAI_API_KEY')
MODEL = "gpt-4o"

def validate_image_data(image_data):
    """Ensure image data is properly formatted"""
    if not image_data:
        return False
    if isinstance(image_data, str):
        if image_data.startswith('data:image'):
            return True
        try:
            base64.b64decode(image_data)
            return True
        except:
            return False
    return False

def extract_image_bytes(image_data):
    """Extract raw image bytes from different formats"""
    if image_data.startswith('data:image'):
        return base64.b64decode(image_data.split(',')[1])
    return base64.b64decode(image_data)

def analyze_image(image_path):
    """Get reconstruction instructions with enhanced prompting"""
    with open(image_path, "rb") as img_file:
        response = openai.chat.completions.create(
            model=MODEL,
            messages=[
                {
                    "role": "system",
                    "content": (
                        "You are a PowerPoint reconstruction AI. Generate COMPLETE Python code to recreate ALL elements from the screenshot.\n\n"
                        "STRICT REQUIREMENTS:\n"
                        "1. Output this function structure:\n"
                        "def create_slide():\n"
                        "    prs = Presentation()\n"
                        "    slide = prs.slides.add_slide(prs.slide_layouts[6])\n"
                        "    # Add elements here\n"
                        "    return prs\n\n"
                        "2. For ALL visual elements:\n"
                        "   - Text boxes: Preserve exact content, formatting, and position\n"
                        "   - Shapes: Recreate rectangles, lines, arrows using add_shape()\n"
                        "   - Position: Use Inches() with exact measurements\n"
                        "   - Styling: Match colors (fill/line), line weights, and dash styles\n\n"
                        "3. Shape handling rules:\n"
                        "   - Rectangles: MSO_SHAPE.RECTANGLE\n"
                        "   - Lines/Arrows: MSO_SHAPE.STRAIGHT_CONNECTOR\n"
                        "   - Ovals: MSO_SHAPE.OVAL\n"
                        "   - Set fill: shape.fill.solid(), .fore_color.rgb\n"
                        "   - Set line: shape.line.color.rgb, .width\n\n"
                        "4. Prohibited:\n"
                        "   - Omitted elements\n"
                        "   - Approximate positioning\n"
                        "   - Color approximations\n"
                        "   - DO NOT use .add_freeform() – it is NOT supported in python-pptx\n"
                        "   - Only use add_shape(), add_textbox(), add_picture(), add_table(), add_chart()\n"
                    )
                },
                {
                    "role": "user",
                    "content": [
                        {
                            "type": "text", 
                            "text": "Rebuild this screenshot including ALL shapes and text elements with exact positioning:"
                        },
                        {
                            "type": "image_url",
                            "image_url": {
                                "url": f"data:image/png;base64,{base64.b64encode(img_file.read()).decode('utf-8')}"
                            }
                        }
                    ]
                }
            ],
            max_tokens=4000,
            temperature=0.1,
            top_p=0.9
        )
    return response.choices[0].message.content

@app.route('/')
def home():
    return send_from_directory('static', 'index.html')

@app.route('/convert', methods=['POST'])
def convert():
    try:
        if not request.is_json:
            return jsonify({"error": "Request must be JSON"}), 400
            
        data = request.get_json()
        if not data or 'image' not in data:
            return jsonify({"error": "No image provided"}), 400
        
        if not validate_image_data(data['image']):
            return jsonify({"error": "Invalid image format"}), 400
            
        img_bytes = extract_image_bytes(data['image'])
        
        try:
            Image.open(io.BytesIO(img_bytes)).verify()
        except:
            return jsonify({"error": "Invalid image data"}), 400

        with tempfile.NamedTemporaryFile(delete=False, suffix='.png') as tmp:
            tmp.write(img_bytes)
            img_path = tmp.name

        analysis = analyze_image(img_path)
        if not analysis:
            return jsonify({"error": "AI analysis failed"}), 400

        code_match = re.search(r'```python(.*?)```', analysis, re.DOTALL)
        if not code_match:
            return jsonify({"error": "No valid Python code generated"}), 400
            
        ppt_code = code_match.group(1).strip()
        
        if 'add_freeform' in ppt_code:
            return jsonify({"error": "Generated code uses unsupported method: add_freeform()"}), 400

        # Add shape-related imports if missing
        required_imports = [
            'MSO_SHAPE', 'MSO_CONNECTOR', 'MSO_LINE', 'MSO_FILL'
        ]
        if not all(imp in ppt_code for imp in required_imports):
            ppt_code = (
                "from pptx.enum.shapes import MSO_SHAPE, MSO_CONNECTOR\n"
                "from pptx.enum.dml import MSO_LINE, MSO_FILL\n" + ppt_code
            )

        exec_globals = {
            'Presentation': Presentation,
            'Inches': Inches,
            'Pt': Pt,
            'RGBColor': RGBColor,
            'PP_ALIGN': PP_ALIGN,
            'MSO_SHAPE': MSO_SHAPE,
            'MSO_CONNECTOR': MSO_CONNECTOR,
            'MSO_LINE': MSO_LINE,
            'MSO_FILL': MSO_FILL
        }

        try:
            exec(ppt_code, exec_globals)
            if 'create_slide' not in exec_globals:
                return jsonify({"error": "Missing create_slide function"}), 400
                
            prs = exec_globals['create_slide']()
            
            if not hasattr(prs, 'save'):
                return jsonify({"error": "Invalid presentation object"}), 400

            ppt_stream = io.BytesIO()
            prs.save(ppt_stream)
            ppt_stream.seek(0)

            return send_file(
                ppt_stream,
                as_attachment=True,
                download_name='output.pptx',
                mimetype='application/vnd.openxmlformats-officedocument.presentationml.presentation'
            )

        except Exception as e:
            return jsonify({"error": f"Execution error: {str(e)}"}), 400

    except Exception as e:
        return jsonify({"error": f"Server error: {str(e)}"}), 500
    finally:
        if 'img_path' in locals() and os.path.exists(img_path):
            try:
                os.unlink(img_path)
            except:
                pass

if __name__ == '__main__':
    app.run(host='0.0.0.0', port=5000, debug=True)
