# pptx_service.py: PPT creation

from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE, MSO_CONNECTOR
from pptx.enum.dml import MSO_LINE, MSO_FILL
import re


def create_ppt_from_structure(title, bullets):
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = title
    content = slide.placeholders[1]
    content.text = "\n".join(bullets)
    return prs


def create_ppt_from_gpt_code(code):
    match = re.search(r'```python(.*?)```', code, re.DOTALL)
    if not match:
        raise ValueError("no valid pyhton code")

    ppt_code = match.group(1).strip()
    if 'add_freeform' in ppt_code:
        raise ValueError("forbidden method: add_freeform()")

    exec_globals = {
        'Presentation': Presentation,
        'Inches': Inches,
        'Pt': Pt,
        'RGBColor': RGBColor,
        'PP_ALIGN': PP_ALIGN,
        'MSO_SHAPE': MSO_SHAPE,
        'MSO_CONNECTOR': MSO_CONNECTOR,
        'MSO_LINE': MSO_LINE,
        'MSO_FILL': MSO_FILL
    }

    exec(ppt_code, exec_globals)
    if 'create_slide' not in exec_globals:
        raise ValueError("create_slide() not defined")

    return exec_globals['create_slide']()

from pptx import Presentation
from pptx.util import Inches, Pt

def create_ppt_from_elements(elements):
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    for el in elements:
        textbox = slide.shapes.add_textbox(
            Inches(el['left']),
            Inches(el['top']),
            Inches(el['width']),
            Inches(el['height'])
        )
        frame = textbox.text_frame
        p = frame.paragraphs[0]
        run = p.add_run()
        run.text = el['text']
        run.font.size = Pt(el['font_size'])

    return prs
