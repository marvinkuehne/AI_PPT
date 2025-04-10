import os
import base64
import io
import tempfile
import re
from flask import Flask, request, jsonify, send_file, send_from_directory
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE, MSO_CONNECTOR
from pptx.enum.dml import MSO_LINE, MSO_FILL
from PIL import Image, ImageEnhance, ImageFilter
import openai
from flask_cors import CORS

app = Flask(__name__, static_folder='static')
CORS(app, origins=["*"]) 

# Configuration
openai.api_key = os.getenv('OPENAI_API_KEY')
MODEL = "gpt-4o"

def validate_image_data(image_data):
    """Ensure image data is properly formatted"""
    if not image_data:
        return False
    if isinstance(image_data, str):
        if image_data.startswith('data:image'):
            return True
        try:
            base64.b64decode(image_data)
            return True
        except:
            return False
    return False

def extract_image_bytes(image_data):
    """Extract raw image bytes from different formats"""
    if image_data.startswith('data:image'):
        return base64.b64decode(image_data.split(',')[1])
    return base64.b64decode(image_data)

def analyze_image(image_path):
    """Send image to OpenAI and get slide generation code."""
    with open(image_path, "rb") as img_file:
        image_data = base64.b64encode(img_file.read()).decode('utf-8')
        response = openai.chat.completions.create(
            model=MODEL,
            messages=[
                {
                    "role": "system",
                    "content": (
                        "You are a skilled Python developer converting diagrams from sketches into PowerPoint slides using python-pptx.\n\n"
                        "Guidelines:\n"
                        "1. Output ONLY a function:\n"
                        "```python\n"
                        "def create_slide():\n"
                        "    prs = Presentation()\n"
                        "    slide = prs.slides.add_slide(prs.slide_layouts[6])\n"
                        "    # Add elements\n"
                        "    return prs\n"
                        "```\n\n"
                        "2.Given an image (like a flowchart or layout), your job is to:\n"
                        "- Place them on a PowerPoint slide using absolute coordinates (pixels or points).\n"
                        "- Use `Inches(x), Inches(y)` for position and size when using python-pptx.+n"
                        " Use `Inches(x), Inches(y)` for position and size when using python-pptx.\n"
                        "- **Ensure all content fits within the standard slide size of 10 x 7.5 inches.** If needed, proportionally scale or pad layout so that nothing is clipped.\n"
                        "- For each element, use the closest matching `MSO_AUTO_SHAPE_TYPE` (e.g., RECTANGLE, OVAL, DIAMOND).\n"
                        "- Use `slide.shapes.add_shape(...)` for shapes.\n"
                        "- Use `slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, ...)` for arrows or lines between shapes. You may call `.BeginConnect()` and `.EndConnect()` if attaching to shapes.\n"
                        "- **Do NOT use `MSO_AUTO_SHAPE_TYPE.LINE` – it doesn’t exist and will cause an error.**\n"
                        "- Use `Inches` consistently for positioning.\n"
                        "- Add minimal placeholder text to shapes (e.g., Step 1, Decision, etc.).\n"
                        "- Don't add background, images, or titles unless visible in the sketch.\n"
                        "- Keep your response focused only on the final code inside the `generate_ppt_from_sketch` function (no explanations or markdown).\n"
                        "3. No add_freeform or placeholders.\n"
                        "4. Rebuild all visible elements:\n"
                        "- Exact text content, position, font size, alignment, bold/italic\n"
                        "- Shapes: rectangles, arrows, ovals (use MSO_SHAPE.RECTANGLE, MSO_CONNECTOR.STRAIGHT, etc.)\n"
                        "- Use Inches(x, y) for positioning\n"
                        "- Use shape.fill.fore_color.rgb, shape.line.color.rgb, shape.line.width, shape.line.dash_style\n"
                        "- Align elements, avoid crowding, ensure spacing and symmetry\n"
                        "- Treat rough or messy inputs as inspiration—rebuild cleanly and professionally\n\n"
                        "5. Make sure all elements fit within the standard PowerPoint slide size (10 inches × 7.5 inches). Nothing should go beyond the slide bounds.\n\n"
                        "Restrictions:\n"
                        "- No curved arrows or unsupported methods\n"
                        "- Return only the function code inside a single code block—no comments or explanations\n"
                        "- Do not skip any visible content\n"
                        "- Keep code compact, readable, and complete\n\n"
                        "Important:\n"
                        "You must ensure that no shape, connector, or label goes outside the visible PowerPoint slide frame.\n"
                        "If spacing is too tight, proportionally shrink all positions and sizes to fit within slide bounds.\n"
                    )
                },
                {
                    "role": "user",
                    "content": [
                        {
                            "type": "text", 
                            "text": "Rebuild this screenshot including ALL shapes and text elements with exact positioning:"
                        },
                        {
                            "type": "image_url",
                            "image_url": {
                                "url": f"data:image/png;base64,{image_data}"
                            }
                        }
                    ]
                }
            ],
            max_tokens=4000,
            temperature=0.1,
            top_p=0.95
        )
    return response.choices[0].message.content


@app.route('/convert', methods=['POST'])
def convert():
    print(openai.api_key)
    try:
        if not request.is_json:
            return jsonify({"error": "Request must be JSON"}), 400
            
        data = request.get_json()
        if not data or 'image' not in data:
            return jsonify({"error": "No image provided"}), 400
        
        if not validate_image_data(data['image']):
            return jsonify({"error": "Invalid image format"}), 400
            
        img_bytes = extract_image_bytes(data['image'])
        
        try:
            Image.open(io.BytesIO(img_bytes)).verify()
        except:
            return jsonify({"error": "Invalid image data"}), 400

        with tempfile.NamedTemporaryFile(delete=False, suffix='.png') as tmp:
            tmp.write(img_bytes)
            img_path = tmp.name

        analysis = analyze_image(img_path)
        if not analysis:
            return jsonify({"error": "AI analysis failed"}), 400
        
        # Log the full AI response for debugging
        print("AI Analysis Response:", analysis)
        
        code_match = re.search(r'```python(.*?)```', analysis, re.DOTALL)
        if not code_match:
            return jsonify({"error": "No valid Python code generated"}), 400
            
        ppt_code = code_match.group(1).strip()
        
        
        if 'add_freeform' in ppt_code:
            return jsonify({"error": "Generated code uses unsupported method: add_freeform()"}), 400
        
        # Second-pass refinement
        try:
            refine_response = openai.chat.completions.create(
                model=MODEL,
                messages=[
                    {
                        "role": "system",
                        "content": "You are a professional PowerPoint code optimizer. Refactor and polish this Python code that generates slides using python-pptx. Improve layout alignment, spacing, visual polish. Do not change the content or logic, only improve aesthetics. Output ONLY the new function code inside a code block."
                    },
                    {
                        "role": "user",
                        "content": f"Improve this PowerPoint layout:\n```python\n{ppt_code}\n```"
                    }
                ],
                temperature=0.2,
                max_tokens=3000
            )
            refined_match = re.search(r'```python(.*?)```', refine_response.choices[0].message.content, re.DOTALL)
            if refined_match:
                ppt_code = refined_match.group(1).strip()
                print("🔁 Refined Code:\n", ppt_code)
            else:
                print("⚠️ Refinement failed, using original code.")

        except Exception as e:
            print("⚠️ Refinement error:", str(e))


        # Add shape-related imports if missing
        required_imports = [
            'MSO_SHAPE', 'MSO_CONNECTOR', 'MSO_LINE', 'MSO_FILL'
        ]
        if not all(imp in ppt_code for imp in required_imports):
            ppt_code = (
                "from pptx.enum.shapes import MSO_SHAPE, MSO_CONNECTOR\n"
                "from pptx.enum.dml import MSO_LINE, MSO_FILL\n" + ppt_code
            )

        # Correct known invalid attributes
        ppt_code = ppt_code.replace('MSO_AUTO_SHAPE_TYPE.STRAIGHT_CONNECTOR', 'MSO_CONNECTOR.STRAIGHT')
        ppt_code = ppt_code.replace('MSO_AUTO_SHAPE_TYPE.DASH', 'MSO_LINE.DASH')

        # Additional validation to ensure no unsupported attributes are present
        unsupported_attributes = ['MSO_AUTO_SHAPE_TYPE.STRAIGHT_CONNECTOR', 'MSO_AUTO_SHAPE_TYPE.DASH']
        for attr in unsupported_attributes:
            if attr in ppt_code:
                return jsonify({"error": f"Generated code uses unsupported attribute: {attr}"}), 400

        exec_globals = {
            'Presentation': Presentation,
            'Inches': Inches,
            'Pt': Pt,
            'RGBColor': RGBColor,
            'PP_ALIGN': PP_ALIGN,
            'MSO_SHAPE': MSO_SHAPE,
            'MSO_CONNECTOR': MSO_CONNECTOR,
            'MSO_LINE': MSO_LINE,
            'MSO_FILL': MSO_FILL
        }

        try:
            exec(ppt_code, exec_globals)
            if 'create_slide' not in exec_globals:
                return jsonify({"error": "Missing create_slide function"}), 400
            
            prs = exec_globals['create_slide']()
        
            if not hasattr(prs, 'save'):
                return jsonify({"error": "Invalid presentation object"}), 400

            ppt_stream = io.BytesIO()
            prs.save(ppt_stream)
            ppt_stream.seek(0)

            return send_file(
                ppt_stream,
                as_attachment=True,
                download_name='output.pptx',
                mimetype='application/vnd.openxmlformats-officedocument.presentationml.presentation'
            )

        except Exception as e:
            # Log the exception for debugging
            print("Execution error:", str(e))
            return jsonify({"error": f"Execution error: {str(e)}"}), 400

    except Exception as e:
        # Log the exception for debugging
        print("Server error:", str(e))
        return jsonify({"error": f"Server error: {str(e)}"}), 500
    finally:
        if 'img_path' in locals() and os.path.exists(img_path):
            try:
                os.unlink(img_path)
            except:
                pass

if __name__ == '__main__':
    app.run(host='0.0.0.0', port=5000, debug=True)